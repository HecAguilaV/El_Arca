import os
import hashlib
from pathlib import Path
from pypdf import PdfReader
from docx import Document
from pptx import Presentation
import logging

# Configurar logging
logging.basicConfig(level=logging.INFO)
logger = logging.getLogger("ArcaCore")

class ArcaScanner:
    def __init__(self, root_path):
        self.root_path = Path(root_path)
        self.files_index = []

    def scan_directory(self):
        """Escanea recursivamente el directorio y retorna metadatos básicos."""
        results = []
        for entry in self.root_path.rglob("*"):
            if entry.is_file():
                # Ignorar archivos ocultos o del sistema
                if entry.name.startswith("."):
                    continue
                
                stats = entry.stat()
                results.append({
                    "path": str(entry.relative_to(self.root_path)),
                    "full_path": str(entry),
                    "size": stats.st_size,
                    "extension": entry.suffix.lower(),
                    "filename": entry.name
                })
        self.files_index = results
        return results

class FileProcessor:
    @staticmethod
    def get_file_hash(filepath):
        """Calcula hash SHA256 para detectar duplicados exactos."""
        hasher = hashlib.sha256()
        with open(filepath, 'rb') as f:
            buf = f.read(65536)
            while len(buf) > 0:
                hasher.update(buf)
                buf = f.read(65536)
        return hasher.hexdigest()

    @staticmethod
    def extract_text(filepath):
        path = Path(filepath)
        ext = path.suffix.lower()
        text = ""

        try:
            if ext == ".pdf":
                reader = PdfReader(filepath)
                for page in reader.pages:
                    text += page.extract_text() + "\n"
            
            elif ext == ".docx":
                doc = Document(filepath)
                for para in doc.paragraphs:
                    text += para.text + "\n"
            
            elif ext in [".pptx", ".ppt"]:
                # Nota: .ppt antiguo no soportado por python-pptx, solo .pptx
                if ext == ".pptx":
                    prs = Presentation(filepath)
                    for slide in prs.slides:
                        for shape in slide.shapes:
                            if hasattr(shape, "text"):
                                text += shape.text + "\n"
            
            elif ext == ".txt":
                with open(filepath, "r", encoding="utf-8", errors="ignore") as f:
                    text = f.read()

        except Exception as e:
            logger.error(f"Error procesando {filepath}: {e}")
            return f"[Error de Lectura: {str(e)}]"

        return text.strip()

# Instancia global para pruebas simples
if __name__ == "__main__":
    # Test
    pass
