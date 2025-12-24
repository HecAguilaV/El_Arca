import os
import hashlib
import json
import csv
import logging
import argparse
import shutil
from pathlib import Path
from typing import List, Dict, Any
import fitz  # PyMuPDF
from tqdm import tqdm

# Intentamos importar pptx, si falla, no se detiene el script
try:
    from pptx import Presentation
    HAS_PPTX = True
except ImportError:
    HAS_PPTX = False

# Configuración de Logging
logging.basicConfig(level=logging.INFO, format='%(asctime)s - %(levelname)s - %(message)s')
logger = logging.getLogger(__name__)

# --- CLASIFICACIÓN INTELIGENTE ---
KEYWORDS = {
    "Escuela Dominical": ["manualidad", "colorear", "niños", "infantil", "dinámica", "recortar", "títeres", "clase bíblica"],
    "Seminario": ["teología", "doctrina", "exégesis", "hermenéutica", "sistemática", "historia", "griego", "hebreo", "comentario"],
    "Multimedia": ["afiche", "banner", "logo", "presentación", "diapositiva"],
    "Música": ["himnario", "partitura", "acordes", "coros"],
}

class ArkInventory:
    def __init__(self, root_dir: str):
        self.root_dir = Path(root_dir)
        self.inventory = []
        self.hashes = set()
        self.isbn_registry = {}  # Futuro: cache de ISBNs

    def calculate_md5(self, file_path: Path, chunk_size: int = 4096) -> str:
        """Calcula hash MD5."""
        hash_md5 = hashlib.md5()
        try:
            with open(file_path, "rb") as f:
                for chunk in iter(lambda: f.read(chunk_size), b""):
                    hash_md5.update(chunk)
            return hash_md5.hexdigest()
        except Exception as e:
            logger.error(f"Error hash {file_path}: {e}")
            return ""

    def convert_epub_to_pdf(self, epub_path: Path) -> Path:
        """Convierte EPUB a PDF usando PyMuPDF."""
        try:
            pdf_path = epub_path.with_suffix(".pdf")
            if pdf_path.exists():
                return pdf_path  # Ya existe, no reconvertir

            doc = fitz.open(epub_path)
            pdf_bytes = doc.convert_to_pdf()
            with open(pdf_path, "wb") as f:
                f.write(pdf_bytes)
            doc.close()
            logger.info(f"Convertido: {epub_path.name} -> PDF")
            return pdf_path
        except Exception as e:
            logger.error(f"Error convirtiendo EPUB {epub_path}: {e}")
            return None

    def get_metadata_and_text(self, file_path: Path) -> Dict[str, Any]:
        """Extrae metadatos y texto según el formato."""
        meta = {
            "page_count": 0,
            "text_preview": "",
            "is_encrypted": False,
            "error": None,
            "format": file_path.suffix.lower().replace(".", "")
        }
        
        try:
            # --- PDF ---
            if file_path.suffix.lower() == ".pdf":
                doc = fitz.open(file_path)
                meta["page_count"] = doc.page_count
                meta["is_encrypted"] = doc.is_encrypted
                text = []
                for i in range(min(5, doc.page_count)):
                    try:
                        text.append(doc.load_page(i).get_text())
                    except: pass
                meta["text_preview"] = " ".join(text)
                doc.close()

            # --- PPTX ---
            elif file_path.suffix.lower() in [".pptx", ".ppt"] and HAS_PPTX:
                prs = Presentation(file_path)
                meta["page_count"] = len(prs.slides)
                text = []
                for i, slide in enumerate(prs.slides):
                    if i >= 5: break
                    for shape in slide.shapes:
                        if hasattr(shape, "text"):
                            text.append(shape.text)
                meta["text_preview"] = " ".join(text)

        except Exception as e:
            meta["error"] = str(e)
            
        return meta

    def classify_content(self, filename: str, text: str, size_bytes: int) -> Dict[str, str]:
        """Determina Categoría y Etiquetas."""
        content_lower = (filename + " " + text).lower()
        
        # 1. Categoría por Palabras Clave
        category = "General"
        tags = []
        
        for cat, kw_list in KEYWORDS.items():
            for kw in kw_list:
                if kw in content_lower:
                    tags.append(kw)
                    if category == "General":
                        category = cat # Asigna la primera categoría fuerte encontrada

        # 2. Refinamiento por extensión
        if "ppt" in filename.lower():
            category = "Multimedia"
        
        # 3. Refinamiento por tamaño (si sigue siendo General)
        if category == "General":
            if size_bytes > 2 * 1024 * 1024:
                category = "Libro/Artículo"
            else:
                category = "Boletín/Breve"

        return {"category": category, "tags":  list(set(tags))} # Unique tags

    def scan(self):
        """Escaneo principal recursivo."""
        logger.info(f"Escaneando: {self.root_dir}")
        
        # Extensiones soportadas
        extensions = ["*.pdf", "*.epub", "*.pptx", "*.ppt"]
        all_files = []
        for ext in extensions:
            all_files.extend(list(self.root_dir.rglob(ext)))

        for file_path in tqdm(all_files, desc="Procesando", unit="file"):
            try:
                # Conversión al vuelo de EPUB
                current_path = file_path
                if file_path.suffix.lower() == ".epub":
                    converted = self.convert_epub_to_pdf(file_path)
                    if converted:
                        current_path = converted
                    else:
                        continue # Skip si falló conversión

                # Datos básicos
                stats = current_path.stat()
                size_bytes = stats.st_size
                
                # Deduplicación
                file_hash = self.calculate_md5(current_path)
                is_duplicate = file_hash in self.hashes
                if file_hash: self.hashes.add(file_hash)

                # Metadatos y Contenido
                meta = self.get_metadata_and_text(current_path)
                
                # Clasificación
                classification = self.classify_content(current_path.name, meta["text_preview"], size_bytes)

                entry = {
                    "path": str(current_path.relative_to(self.root_dir)),
                    "filename": current_path.name,
                    "format": meta["format"],
                    "size_bytes": size_bytes,
                    "category": classification["category"],
                    "tags": ", ".join(classification["tags"]),
                    "md5_hash": file_hash,
                    "is_duplicate": is_duplicate,
                    "page_count": meta["page_count"],
                    "error": meta["error"]
                }
                
                self.inventory.append(entry)

            except Exception as e:
                logger.error(f"Fallo en {file_path.name}: {e}")

    def save_results(self):
        json_path = self.root_dir / "inventario_arca.json"
        with open(json_path, 'w', encoding='utf-8') as f:
            json.dump(self.inventory, f, indent=4, ensure_ascii=False)
        logger.info(f"Guardado JSON: {json_path}")

def main():
    parser = argparse.ArgumentParser()
    # Ajuste: La raíz por defecto ahora es relativa a 'scripts/' hacia 'public/library'
    default_root = Path(__file__).parent.parent / "public" / "library"
    parser.add_argument("--root", type=str, default=str(default_root), help="Directorio raíz")
    args = parser.parse_args()

    inv = ArkInventory(args.root)
    inv.scan()
    
    # Mover el JSON generado al lugar correcto (public/data.json)
    # save_results guarda en root_dir (public/library/inventario_arca.json)
    inv.save_results()

    source_json = Path(args.root) / "inventario_arca.json"
    dest_json = Path(args.root).parent / "data.json"
    
    if source_json.exists():
        shutil.move(str(source_json), str(dest_json))
        logger.info(f"✅ Datos publicados en: {dest_json}")

if __name__ == "__main__":
    main()
